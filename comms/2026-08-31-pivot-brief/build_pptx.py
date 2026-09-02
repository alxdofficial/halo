"""Build the task/workflow/data briefing, in the established house style.

Scope for this revision, per Alex: only (1) the three real-world scenarios and the
task names, (2) per-task explanation + baseline-encoder -> learnable-head ->
task-head -> prediction workflow, (3) how train/eval data is constructed, and
(4) one detailed per-dataset numbers table. Everything else from the earlier
pivot narrative (evidence, status, decisions) is dropped.

Same geometry and palette as prior decks: 13.333x7.5in, Arial, title 25pt bold
#1A1A1A at (0.60, 0.28), italic deck line 12pt #34618E, card panels
#F4F6F9/#CCCCCC, cream takeaway banner #FBF3DF, tables with an #EFEFEF header.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path("comms/2026-08-31-pivot-brief")

INK, MUTED, BLUE, GREEN = "1A1A1A", "666666", "34618E", "3F6B4A"
CREAM, CARD, LINE, HDR, HLROW = "FBF3DF", "F4F6F9", "CCCCCC", "EFEFEF", "EFF4EE"
TASK = {1: "2A64A8", 2: "7B4B94", 3: "B5651D"}
FONT = "Arial"


def rgb(h):
    return RGBColor.from_string(h)


prs = Presentation()
prs.slide_width, prs.slide_height = Emu(12191695), Emu(6858000)
BLANK = prs.slide_layouts[6]


def new():
    return prs.slides.add_slide(BLANK)


def textbox(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def write(tf, lines, size, colour=INK, bold=False, italic=False, space=6, first=True):
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if (first and i == 0) else tf.add_paragraph()
        p.space_after = Pt(space)
        for j, chunk in enumerate(str(line).split("**")):
            if not chunk:
                continue
            r = p.add_run()
            r.text = chunk
            r.font.size, r.font.name = Pt(size), FONT
            r.font.bold = bold or (j % 2 == 1)
            r.font.italic = italic
            r.font.color.rgb = rgb(colour)
    return tf


def title(s, text):
    write(textbox(s, 0.60, 0.28, 12.10, 0.55), text, 25, INK, bold=True)


def deckline(s, text):
    write(textbox(s, 0.60, 0.83, 12.10, 0.34), text, 12, BLUE, italic=True)


def panel(s, l, t, w, h, fill=CARD, line=LINE, lw=9525):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    sh.line.color.rgb = rgb(line); sh.line.width = lw
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def banner(s, text, top=6.15):
    panel(s, 0.60, top, 12.13, 0.66, fill=CREAM, lw=6350)
    tf = textbox(s, 0.80, top, 11.73, 0.66, anchor=MSO_ANCHOR.MIDDLE)
    write(tf, text, 13.5, INK, bold=True, space=0)


def colhead(s, l, t, w, text):
    write(textbox(s, l, t, w, 0.32), text, 15, INK, bold=True)


def bullets(s, l, t, w, h, lines, size=13.0):
    write(textbox(s, l, t, w, h), lines, size, INK, space=8)


def table(s, rows, l, t, widths, row_h=0.34, size=10.0, hi_rows=(), left_cols=(0,)):
    shape = s.shapes.add_table(len(rows), len(rows[0]), Inches(l), Inches(t),
                               Inches(sum(widths)), Inches(row_h * len(rows)))
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    for i, wdt in enumerate(widths):
        tbl.columns[i].width = Inches(wdt)
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(row_h)
        for ci, cell_spec in enumerate(row):
            text, bold = (cell_spec if isinstance(cell_spec, tuple) else (cell_spec, False))
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill = HDR if ri == 0 else (HLROW if ri in hi_rows else "FFFFFF")
            cell.fill.solid(); cell.fill.fore_color.rgb = rgb(fill)
            colour = GREEN if (ri in hi_rows and ri != 0) else INK
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci in left_cols else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(text)
            r.font.size, r.font.name = Pt(size), FONT
            r.font.bold = bold or ri == 0
            r.font.color.rgb = rgb(colour)
    return shape


def picture(s, path, l, t, w, h):
    """Fit an image into the (l, t, w, h) box, preserving aspect, centred."""
    from PIL import Image
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    s.shapes.add_picture(str(path), Inches(l + (w - dw) / 2), Inches(t + (h - dh) / 2),
                         Inches(dw), Inches(dh))


def stripe(s, l, t, h, color, w=0.09):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(color)
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


# ================================================================ 1  title
s = new()
write(textbox(s, 0.60, 2.20, 12.00, 0.90),
      "Movement Monitoring — Tasks, Workflows, Data", 32, INK, bold=True)
write(textbox(s, 0.60, 3.15, 12.00, 0.50),
      "1 Sep 2026 · branch application-motion-monitoring", 16, MUTED)
write(textbox(s, 0.60, 3.85, 12.00, 0.60),
      "Three real-world scenarios, three tasks  ·  encoder-to-prediction workflow  ·  "
      "how the training and evaluation data is built", 14.5, INK)

# ================================================================ 2  three scenarios / three tasks
s = new()
title(s, "Three real-world scenarios, three tasks")
deckline(s, "Each task imitates one concrete situation a clinician, patient, or ergonomist is in.")
cards = [
    (1, "DETECT",
     "Scenario: a clinician demonstrates an exercise; later, did it happen again?",
     "Given a few reference executions, locate that same movement in a continuous "
     "recording — different session, speed, or device mounting."),
    (2, "COMPARE",
     "Scenario: is this week's exercise different from last week's, for this person?",
     "For repeated executions of one task, quantify how it changed relative to the "
     "person's own accepted baseline — not a population norm."),
    (3, "DISCOVER",
     "Scenario: an ergonomist reviews an unlabeled shift and finds what repeats.",
     "Cluster recurring motion motifs in a continuous recording with no reference "
     "and no predefined activity vocabulary."),
]
for i, (tid, head, scenario, desc) in enumerate(cards):
    x = 0.60 + i * 4.12
    panel(s, x, 1.55, 3.85, 4.60)
    stripe(s, x, 1.55, 4.60, TASK[tid])
    picture(s, ROOT / f"figures/icon_task{tid}.png", x + 0.24, 1.75, 0.85, 0.85)
    write(textbox(s, x + 1.28, 1.90, 2.35, 0.55), f"TASK {tid}", 12, TASK[tid], bold=True, space=0)
    write(textbox(s, x + 1.28, 2.22, 2.35, 0.55), head, 17, INK, bold=True, space=0)
    write(textbox(s, x + 0.28, 2.90, 3.35, 1.15), scenario, 12.5, INK, bold=True, space=0)
    write(textbox(s, x + 0.28, 4.05, 3.35, 1.95), desc, 12, INK, space=0)
banner(s, "One shared primitive underneath all three: timestamped patch embeddings, never "
          "pooled into a single vector per recording.", top=6.35)

# ================================================================ 3-4  Task 1
s = new()
stripe(s, 0.0, 0.0, 7.5, TASK[1])
title(s, "Task 1 — Detect")
deckline(s, "Given one or a few reference executions, did the movement occur in a later "
            "recording, and where?")
colhead(s, 0.60, 1.55, 5.85, "Real-world scenario")
bullets(s, 0.60, 1.97, 5.85, 3.50, [
    "A clinician records several acceptable repetitions of a prescribed exercise in "
    "clinic — or a worker demonstrates an arbitrary personal task, “exercise one.”",
    "The system later scans a continuous home or field recording, counts detected "
    "repetitions, and returns the matching intervals for inspection.",
    "The reference's language label is irrelevant — names are never a required input.",
], 12.5)
colhead(s, 6.85, 1.55, 5.85, "What the system must handle")
bullets(s, 6.85, 1.97, 5.85, 3.50, [
    "The later execution may differ in speed, amplitude, ordinary repetition "
    "variability, session, or device mounting.",
    "Reference and query must be independent executions — never two windows from "
    "one bout.",
    "Recordings with no target occurrence, and hard negatives from the same person, "
    "device, and activity family.",
], 12.5)
banner(s, "Output: event intervals, match scores, and the retrieved reference evidence "
          "behind each match.", top=6.15)

s = new()
stripe(s, 0.0, 0.0, 7.5, TASK[1])
title(s, "Task 1 — workflow")
deckline(s, "Reference executions and a continuous query timeline, both through the same "
            "frozen encoder.")
picture(s, ROOT / "figures/workflow_task1.png", 0.60, 1.72, 12.13, 3.35)
banner(s, "The task head is a search over the whole timeline, not a per-window "
          "classifier — background dominates a continuous recording.", top=5.35)

# ================================================================ 5-6  Task 2
s = new()
stripe(s, 0.0, 0.0, 7.5, TASK[2])
title(s, "Task 2 — Compare")
deckline(s, "For executions known to represent the same task, which differences are ordinary "
            "variation and which are a persistent change?")
colhead(s, 0.60, 1.55, 5.85, "Real-world scenario")
bullets(s, 0.60, 1.97, 5.85, 3.50, [
    "Track whether a rehabilitation exercise, or a worker's task execution, is "
    "changing across days or weeks.",
    "The comparison is always against that person's own accepted baseline — not "
    "against another person or a population average.",
    "Alignment should ignore properties like speed while keeping others, like "
    "duration, visible rather than discarding them.",
], 12.5)
colhead(s, 6.85, 1.55, 5.85, "What the system must handle")
bullets(s, 6.85, 1.97, 5.85, 3.50, [
    "Separate test-retest and device-remounting noise from a real persistent "
    "change — the change must exceed that noise floor.",
    "Isolated outliers should be suppressed while persistent drift and its "
    "uncertainty are retained.",
    "Without external ground truth the output is difference, never “quality” or "
    "“improvement.”",
], 12.5)
banner(s, "Output: aligned latent-shape deviation, phase-local deviation, duration / "
          "intensity / smoothness deltas, and uncertainty vs. personal baseline.", top=6.15)

s = new()
stripe(s, 0.0, 0.0, 7.5, TASK[2])
title(s, "Task 2 — workflow")
deckline(s, "Two independent executions of the same declared task, reference and comparison.")
picture(s, ROOT / "figures/workflow_task2.png", 0.60, 1.72, 12.13, 3.35)
banner(s, "The task head aligns by phase, not by clock time, so duration differences "
          "stay visible instead of being warped away.", top=5.35)

# ================================================================ 7-8  Task 3
s = new()
stripe(s, 0.0, 0.0, 7.5, TASK[3])
title(s, "Task 3 — Discover")
deckline(s, "Which coherent motion motifs recur frequently in an unlabeled continuous recording?")
colhead(s, 0.60, 1.55, 5.85, "Real-world scenario")
bullets(s, 0.60, 1.97, 5.85, 3.50, [
    "An ergonomist records a work shift with no reference and no predefined "
    "activity vocabulary.",
    "The system surfaces the most frequently recurring motion motifs for review; "
    "the ergonomist names the ones that matter.",
    "Named motifs are promoted into Task 1 as references for future monitoring.",
], 12.5)
colhead(s, 6.85, 1.55, 5.85, "What the system must handle")
bullets(s, 6.85, 1.97, 5.85, 3.50, [
    "IMU data cannot prove intent — a coherent motif is not assumed meaningful "
    "until a human confirms it.",
    "The same-motion metric trains on arbitrary labeled event identities; "
    "evaluation identities and label names stay hidden from clustering.",
    "Recurrence, stable internal ordering, and within-cluster similarity greater "
    "than similarity to local background.",
], 12.5)
banner(s, "Output: recurring motif clusters — count, duration, cadence, representative "
          "examples, and timeline locations, pending human confirmation.", top=6.15)

s = new()
stripe(s, 0.0, 0.0, 7.5, TASK[3])
title(s, "Task 3 — workflow")
deckline(s, "One continuous, unlabeled recording; no reference, no vocabulary.")
picture(s, ROOT / "figures/workflow_task3.png", 0.60, 1.72, 12.13, 3.35)
banner(s, "The task head never sees a label name — only which candidates were marked "
          "the same event during training.", top=5.35)

# ================================================================ 9  dataset mix
s = new()
title(s, "How the training and evaluation data is built")
deckline(s, "No dataset both trains and evaluates the same arm.")
picture(s, ROOT / "figures/dataset_mix.png", 0.85, 1.35, 11.4, 4.60)
banner(s, "Training sources never contribute to a sealed test, and a sealed source is "
          "never fit or calibrated on.", top=6.15)

# ================================================================ 10  dataset detail table
DATASETS = [
    ("C-MHAD", "240 streams\n(12 subj × 2 apps × 10 runs)", "~119–122 s",
     "8.0 h", "12", "12",
     "5 TV gestures + 7 sit/stand/lie transitions · wrist or waist Shimmer3 IMU 50 Hz + "
     "sync video (unused)"),
    ("WEAR", "24 sessions", "~50 min avg (varies)", "19.0 h", "22", "18 (+NULL)",
     "18 outdoor sports activities · 4 limb Bangle.js watches 50 Hz, 2 arm streams used, "
     "+ egocentric video (unused)"),
    ("MoniPar", "174 weekly sessions", "405–445 s (~7 min)", "~20 h*", "28 (21 PD + 7 ctrl)",
     "9", "9 clinical exercise labels · wrist smartwatch 50 Hz, gravity present, "
     "weekly protocol"),
    ("OCA", "12 sessions (13 parts)", "~14–56 min (varies)", "6.0 h", "5", "6 (+NULL)",
     "6 assembly-phase labels · 4 BNO055 IMUs (2 upper-arm, 2 chest/vest), ~20–27 Hz "
     "native, arm-support flag"),
    ("OpenPack", "416 IMU segments\n(gap-split)", "varies by session", "53.8 h",
     "16 (21 raw IDs, 5 aliased)", "n/a—interval IDs",
     "Nested fine-action / operation / box-cycle identities, not a closed vocabulary · "
     "wrist+arm ATR IMU ~30 Hz, 4 sensors/subject"),
    ("CrossFit", "446 exercise arrays\n(+5,461 rep slices)", "~31 s avg (derived)",
     "3.834 h", "57 codes (50 active;\npaper: 54 people)", "10 (+NULL)",
     "10 exercises, repetition starts cued by watch vibration · wrist+ankle smartwatch "
     "~100 Hz, authors' 10 ms interpolation"),
    ("AIDLAB-HAR", "180 EDF recordings", "short, unstated exact", "not stated",
     "90 codes (not\nverified unique)", "16",
     "13 exercises + 3 background-like classes · chest AIDLAB IMU 50 Hz, acceleration "
     "only (no raw gyro)"),
    ("RecoFit", "126 visits", "varies (continuous)", "79.4 h", "94", "n/a—weak labels",
     "Broad weakly-labeled exercise-set vocabulary + explicit non-exercise background · "
     "right-forearm 50 Hz acceleration + gyroscope"),
]
s = new()
title(s, "Dataset detail")
deckline(s, "The eight newly-verified sources behind Tasks 1–3. Highlighted rows are "
            "sealed eval sources; the rest are train. HARMES (existing corpus source) is "
            "tracked in the corpus docs, not repeated here.")
rows = [["dataset", "recordings", "length / recording", "total length", "subjects",
         "labels (#)", "labels & acquisition"]]
for name, recs, per, total, subj, nlab, desc in DATASETS:
    rows.append([name, recs, per, total, subj, nlab, desc])
table(s, rows, 0.60, 1.55, [1.15, 1.15, 1.35, 0.95, 1.25, 0.95, 5.33],
      row_h=0.56, size=8.3, hi_rows=(1, 2, 3, 4), left_cols=(0, 6))
write(textbox(s, 0.60, 6.66, 12.13, 0.80),
      "* MoniPar total length is derived (sessions × median duration); the release does "
      "not publish a total directly. MoniPar serves Task 1 (cross-week) only — its states "
      "are not independent repetition executions, so it is excluded from Task 2. MoniPar "
      "and HARMES are not yet in the frozen COHORT_V1 manifest (older storage, pending a "
      "reviewed adapter revision) — the other six rows are. “not stated” / “n/a” cells are "
      "genuine gaps in the source documentation, not omissions.",
      9.2, MUTED, italic=True, space=0)

# ---------------------------------------------------------------- geometry check
SLIDE_W, SLIDE_H = 13.3330, 7.5
problems = []
for n, sl in enumerate(prs.slides, start=1):
    boxes = []
    for sh in sl.shapes:
        l, t = Emu(sh.left).inches, Emu(sh.top).inches
        r, b = l + Emu(sh.width).inches, t + Emu(sh.height).inches
        if l < -0.01 or t < -0.01 or r > SLIDE_W + 0.01 or b > SLIDE_H + 0.01:
            problems.append(f"slide {n}: {sh.shape_type} {sh.name} out of bounds "
                            f"({l:.2f},{t:.2f})-({r:.2f},{b:.2f})")
        if sh.shape_type is not None and "PICTURE" in str(sh.shape_type):
            boxes.append(("pic", l, t, r, b))
        elif sh.has_table:
            boxes.append(("table", l, t, r, b))
    banners = [(Emu(sh.top).inches, Emu(sh.top).inches + Emu(sh.height).inches)
               for sh in sl.shapes
               if sh.shape_type is not None and "AUTO_SHAPE" in str(sh.shape_type)
               and getattr(sh.fill, "type", None) == 1
               and str(sh.fill.fore_color.rgb) == CREAM]
    for kind, l, t, r, b in boxes:
        for bt, bb in banners:
            if b > bt + 0.01 and t < bb:
                problems.append(f"slide {n}: {kind} bottom {b:.2f} overlaps banner at {bt:.2f}")
    # A table/picture can silently grow (more rows, a taller image) and collide with a
    # standalone caption or footnote textbox added independently below it.
    text_boxes = [
        (Emu(sh.left).inches, Emu(sh.top).inches,
         Emu(sh.left).inches + Emu(sh.width).inches, Emu(sh.top).inches + Emu(sh.height).inches)
        for sh in sl.shapes
        if sh.has_text_frame and not sh.has_table and sh.text_frame.text.strip()
    ]
    for kind, l, t, r, b in boxes:
        for tl, tt, tr, tb in text_boxes:
            overlap_h = min(r, tr) - max(l, tl)
            overlap_v = min(b, tb) - max(t, tt)
            if overlap_h > 0.01 and overlap_v > 0.01:
                problems.append(
                    f"slide {n}: {kind} ({l:.2f},{t:.2f})-({r:.2f},{b:.2f}) overlaps a "
                    f"text box at ({tl:.2f},{tt:.2f})-({tr:.2f},{tb:.2f})"
                )
print("geometry:", "OK" if not problems else "PROBLEMS")
for line in problems:
    print("  !", line)

out = ROOT / "HALO_Pivot_2026-09-01.pptx"
prs.save(out)
print(f"wrote {out}  ({len(prs.slides._sldIdLst)} slides, {out.stat().st_size/1e6:.2f} MB)")

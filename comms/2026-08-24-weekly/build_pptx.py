"""Build the weekly deck, matching HALO_Weekly_20260811.pptx's house style exactly.

Geometry, type sizes and colours were read off the previous deck: 13.333x7.5in, Arial,
title 25pt bold #1A1A1A at (0.60, 0.28), italic deck line 12pt #34618E at 0.83, card panels
#F4F6F9/#CCCCCC, cream takeaway banner #FBF3DF, tables with an #EFEFEF header and an #EFF4EE
highlight row in #3F6B4A.
"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path("comms/2026-08-24-weekly")
D = json.load(open(ROOT / "data/deck.json"))
KS = ["1", "2", "4", "8", "16"]
NICE = {"HALO": "HALO (ours)", "limubert": "LiMU-BERT", "unimts": "UniMTS", "harnet": "harnet5",
        "crosshar": "CrossHAR", "imagebind": "ImageBind", "normwear": "NormWear"}

INK, MUTED, BLUE, GREEN = "1A1A1A", "666666", "34618E", "3F6B4A"
CREAM, CARD, LINE, HDR, HLROW = "FBF3DF", "F4F6F9", "CCCCCC", "EFEFEF", "EFF4EE"
FONT = "Arial"


def rgb(h):
    return RGBColor.from_string(h)


prs = Presentation()
prs.slide_width, prs.slide_height = Emu(12191695), Emu(6858000)
BLANK = prs.slide_layouts[6]


def new():
    return prs.slides.add_slide(BLANK)


def textbox(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def write(tf, lines, size, colour=INK, bold=False, italic=False, space=6, first=True):
    """lines: str or list. **bold** spans inside a line become bold runs."""
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if (first and i == 0) else tf.add_paragraph()
        p.space_after = Pt(space)
        for j, chunk in enumerate(str(line).split("**")):
            if not chunk:
                continue
            r = p.add_run()
            r.text = chunk
            r.font.size, r.font.name = Pt(size), FONT
            r.font.bold = bold or (j % 2 == 1)
            r.font.italic = italic
            r.font.color.rgb = rgb(colour)
    return tf


def title(s, text):
    write(textbox(s, 0.60, 0.28, 12.10, 0.55), text, 25, INK, bold=True)


def deckline(s, text):
    write(textbox(s, 0.60, 0.83, 12.10, 0.34), text, 12, BLUE, italic=True)


def panel(s, l, t, w, h, fill=CARD, line=LINE, lw=9525):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    sh.line.color.rgb = rgb(line); sh.line.width = lw
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def banner(s, text, top=6.15):
    panel(s, 0.60, top, 12.13, 0.66, fill=CREAM, lw=6350)
    tf = textbox(s, 0.80, top, 11.73, 0.66, anchor=MSO_ANCHOR.MIDDLE)
    write(tf, text, 13.5, INK, bold=True, space=0)


def colhead(s, l, t, w, text):
    write(textbox(s, l, t, w, 0.32), text, 15, INK, bold=True)


def bullets(s, l, t, w, h, lines, size=13.0):
    write(textbox(s, l, t, w, h), lines, size, INK, space=8)


def picture(s, name, l, t, w, h):
    """Fit the image inside the (l, t, w, h) box, preserving aspect, centred.

    Passing width alone lets matplotlib's tight bounding box decide the height, which silently
    pushed three figures underneath their takeaway banner. The box is the contract instead.
    """
    from PIL import Image
    path = ROOT / f"figures/{name}.png"
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    s.shapes.add_picture(str(path), Inches(l + (w - dw) / 2), Inches(t + (h - dh) / 2),
                         Inches(dw), Inches(dh))


def divider(s, eyebrow, head, sub):
    panel(s, 0.0, 2.55, 0.16, 1.90, fill=BLUE, line=BLUE)
    write(textbox(s, 0.62, 2.60, 11.00, 0.36), eyebrow, 13, BLUE, bold=True)
    write(textbox(s, 0.62, 3.02, 11.60, 0.72), head, 32, INK, bold=True)
    write(textbox(s, 0.62, 3.82, 11.60, 0.60), sub, 14, MUTED)


def table(s, rows, l, t, widths, row_h=0.34, size=10.0, hi_rows=(), align_first_left=True):
    shape = s.shapes.add_table(len(rows), len(rows[0]), Inches(l), Inches(t),
                               Inches(sum(widths)), Inches(row_h * len(rows)))
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    for i, wdt in enumerate(widths):
        tbl.columns[i].width = Inches(wdt)
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(row_h)
        for ci, cell_spec in enumerate(row):
            text, bold = (cell_spec if isinstance(cell_spec, tuple) else (cell_spec, False))
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill = HDR if ri == 0 else (HLROW if ri in hi_rows else "FFFFFF")
            cell.fill.solid(); cell.fill.fore_color.rgb = rgb(fill)
            colour = GREEN if (ri in hi_rows and ri != 0) else INK
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if (ci == 0 and align_first_left) else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(text)
            r.font.size, r.font.name = Pt(size), FONT
            r.font.bold = bold or ri == 0
            r.font.color.rgb = rgb(colour)
    return shape


# ---------------------------------------------------------------- data helpers
def headline_rows(regime):
    h = D["headline"]["nearest"]
    order = sorted(h, key=lambda m: -sum(h[m].get(regime, {}).get(k, 0) for k in KS))
    best = {k: max(h[m].get(regime, {}).get(k, 0) for m in h) for k in KS}
    rows = [["model"] + [f"k={k}" for k in KS]]
    hi = []
    for i, m in enumerate(order, start=1):
        cells = [NICE[m]]
        for k in KS:
            v = h[m].get(regime, {}).get(k)
            cells.append((f"{v:.1f}", abs(v - best[k]) < 1e-9))
        rows.append(cells)
        if m == "HALO":
            hi.append(i)
    return rows, tuple(hi)


LADDER = [("pooled_execution_1nn", "HALO encoder + 1-NN"),
          ("support_raw_1nn", "+ keep every recording row"),
          ("full_raw_1nn", "+ add 512-row corpus memory"),
          ("full_reranked_1nn", "+ learned reranker  = native engine")]


def ladder_rows(regime):
    lad = D["ladder"]
    rows = [["what the model does"] + [f"k={k}" for k in KS]]
    for i, (key, label) in enumerate(LADDER):
        cells = [label]
        for k in KS:
            v = lad[key][regime][k]
            if i == 0:
                cells.append(f"{v:.1f}")
            else:
                d = v - lad[LADDER[i - 1][0]][regime][k]
                cells.append(f"{v:.1f}   {d:+.2f}")
        rows.append(cells)
    return rows


# ================================================================ 1  title
s = new()
write(textbox(s, 0.60, 2.20, 12.00, 0.90), "HALO — Weekly Update", 36, INK, bold=True)
write(textbox(s, 0.60, 3.15, 12.00, 0.50),
      "24 Aug 2026 · sealed adaptation_v1 manifest · fingerprint 1bd89d35…", 16, MUTED)
write(textbox(s, 0.60, 3.85, 12.00, 0.40),
      "The best model we have  ·  how everyone is adapted  ·  why our own engine is the problem",
      14.5, INK)

# ================================================================ 2  how it works
s = new()
title(s, "How HALO works")
deckline(s, "One encoder, a frozen text tower, and a memory you drop examples into.")
colhead(s, 0.60, 1.55, 5.85, "The model")
bullets(s, 0.60, 1.97, 5.85, 3.30, [
    "**Input is deliberately unconstrained.** Any channel set, any sampling rate, any body "
    "placement — no resampling to a common grid.",
    "**Filterbank front end.** Fixed Gabor filters at physical frequencies, so 5 Hz means 5 Hz "
    "whether the device sampled at 20 Hz or 200 Hz.",
    "**Temporal trunk.** Three attention layers over one-second patches, width 128 — one 128-d "
    "vector per (patch, sensor) row.",
    "**Text is an input, not a head.** A frozen sentence encoder embeds the label names and a "
    "description of the acquisition setup.",
])
colhead(s, 6.85, 1.55, 5.85, "How it adapts")
bullets(s, 6.85, 1.97, 5.85, 3.30, [
    "**Adaptation is memory, not gradient.** Enrolled examples are encoded once and dropped "
    "into a memory bank. Nothing is retrained, ever.",
    "**The native engine** retrieves the top-64 memory rows for a query, mixes them with the "
    "candidate names, and votes.",
    "**1.01 M trainable parameters** — the smallest model in this comparison by two to three "
    "orders of magnitude.",
])
banner(s, "Today's headline model does not use the engine. It is the HALO encoder read out "
          "with plain nearest-neighbour — slides 11-13 are about why.", top=5.55)

# ================================================================ 3  protocol
s = new()
title(s, "How everything here is measured")
deckline(s, "One sealed manifest, seven held-out datasets, five seeds — identical for every model.")
colhead(s, 0.60, 1.55, 5.85, "The protocol")
bullets(s, 0.60, 1.97, 5.85, 3.30, [
    "**Seven held-out datasets**, in nobody's pretraining: four **ordinary** (everyday activity) "
    "and three **specialized / clinical** (rehab and clinical motion).",
    "**Five seeds**; support and query executions are execution-disjoint and subject-disjoint.",
    "**k = enrolled executions per class.** Each execution is pooled to one vector, so a "
    "56-window recording cannot outweigh a one-window recording.",
    "**Macro F1**, averaged within a dataset, then equally across datasets.",
])
colhead(s, 6.85, 1.55, 5.85, "Two decisions we made this week")
bullets(s, 6.85, 1.97, 5.85, 3.30, [
    "**Three readouts only — 1-NN, prototype, ridge.** All gradient-free. We dropped the fitted "
    "linear head: it needs 200 Adam steps at deployment, which is not the regime we argue for.",
    "**HALO's memory at k ≥ 1 is enrolment only.** No corpus rows, so every model sees exactly "
    "the same evidence. Slide 13 shows this costs us nothing.",
])
banner(s, "Every number in this deck comes from that one manifest. No model was re-tuned for it.",
       top=5.55)

# ================================================================ 4  divider
divider(new(), "RESULTS", "Where the model actually stands",
        "Best checkpoint long_4h_20260821 · frozen encoder + 1-NN · macro F1 on the sealed test set.")

# ================================================================ 5-6 headline
for regime, name, note in [
    ("ordinary", "Ordinary activities",
     "HALO leads at k=1 and is level with LiMU-BERT after that. Call this a tie, not a win."),
    ("specialized_novel", "Specialized and clinical activities",
     "HALO is best at every k, by 6.2 to 9.5 F1 over UniMTS. This is the result worth defending."),
]:
    s = new()
    title(s, name)
    deckline(s, "Every model is a frozen encoder read out with 1-NN. Nothing is fine-tuned.")
    picture(s, f"headline_{regime}", 0.60, 1.42, 7.20, 4.62)
    rows, hi = headline_rows(regime)
    table(s, rows, 7.95, 1.62, [1.55] + [0.64] * 5, row_h=0.36, size=10, hi_rows=hi)
    banner(s, note, top=6.15)

# ================================================================ 7  zero-shot
s = new()
title(s, "Zero-shot, with nothing enrolled")
deckline(s, "No enrolled examples: every model falls back on its own training knowledge, and "
            "chooses among the test dataset's labels only.")
picture(s, "zero_shot", 0.60, 1.35, 12.13, 4.80)
banner(s, "Competitive on everyday activity (37.0, second to CrossHAR); weak on the clinical set "
          "(8.8 against UniMTS at 17.4). One enrolled example takes that 8.8 to 43.2.", top=6.35)

# ================================================================ 8  baseline contracts
s = new()
title(s, "What “adaptation” means for each model at k ≥ 1")
deckline(s, "No backbone is fine-tuned anywhere in this comparison — ours included.")
rows = [["model", "encoder", "params", "feat dim", "pretraining", "input contract", "tuned at k≥1"],
        ["HALO (ours)", "filterbank + 3-layer temporal trunk", "1.01 M", "128",
         "our 18-dataset corpus", "6 ch, native rate", "none"],
        ["LiMU-BERT", "masked IMU encoder", "~62 K", "72",
         "self-pretrained on our corpus", "6 ch @ 20 Hz", "none"],
        ["UniMTS", "graph encoder + CLIP text tower", "~4 M", "512",
         "released ckpt, simulated IMU", "3 ch acc @ 20 Hz", "none"],
        ["harnet5", "ResNet-V2", "~2 M", "512",
         "released ckpt, UK-Biobank", "3 ch acc @ 30 Hz", "none"],
        ["CrossHAR", "masked + contrastive transformer", "~72 K", "72",
         "self-pretrained on our corpus", "6 ch @ 20 Hz", "none"],
        ["NormWear", "channel-independent ViT + frozen LLM", "194 M", "2048",
         "released ckpt, clinical wearables", "6 ch @ 65 Hz", "none"],
        ["ImageBind", "imagebind_huge", "1.2 B", "1024",
         "released ckpt, web-scale", "6 ch @ 200 Hz", "none"]]
table(s, rows, 0.60, 1.50, [1.35, 3.15, 0.85, 0.90, 2.60, 1.85, 1.43], row_h=0.40, size=10,
      hi_rows=(1,))
banner(s, "Every encoder is frozen and no gradient step is taken at adaptation time. The only "
          "thing fitted is the readout — on the k enrolled executions of the target dataset alone.",
       top=5.05)
write(textbox(s, 0.60, 5.85, 12.13, 1.20),
      ["**Why gradient-free is the fair choice.** A fitted head needs a learning rate, a step "
       "count and a weight decay that cannot be right for a 72-d feature and a 2048-d one at "
       "once — any single setting under-tunes somebody, and we would be benchmarking optimisers. "
       "A cosine has nothing to tune and is available to every model, whether its native head "
       "predicts logits (harnet, LiMU-BERT, CrossHAR) or aligns to language (UniMTS, ImageBind, "
       "NormWear)."], 12, INK, space=4)

# ================================================================ 9  readouts
s = new()
title(s, "The three readouts, and what each needs at deployment")
deckline(s, "All three see only the enrolled examples: no corpus data, no other dataset, no query labels.")
rows = [["readout", "parameters fitted", "optimiser", "data it sees"],
        ["1-NN   (primary)", "none", "none — one cosine, one argmax", "the k×C enrolled vectors"],
        ["prototype", "none  (C centroids computed)", "none — mean, then renormalise",
         "the k×C enrolled vectors"],
        ["ridge", "D×C", "one closed-form solve, α fixed at 1", "the k×C enrolled vectors"]]
table(s, rows, 0.60, 1.55, [2.40, 2.95, 3.75, 3.03], row_h=0.42, size=11, hi_rows=(1,))
colhead(s, 0.60, 3.60, 5.85, "Why the memory is enrolment-only")
bullets(s, 0.60, 4.02, 5.85, 1.90, [
    "The baselines' readouts see only enrolled examples. If our engine also carried 512 corpus "
    "rows, k would not mean the same thing across the table.",
    "Measured cost of dropping them: **−0.24 to +0.04 F1**. It buys us nothing, so we give it up "
    "and take the cleaner comparison.",
], 12.0)
colhead(s, 6.85, 3.60, 5.85, "What we still owe")
bullets(s, 6.85, 4.02, 5.85, 1.90, [
    "At k=1, prototype and 1-NN are mathematically identical — one vector per class, and its own "
    "normalised mean is itself.",
    "The text-aligned baselines could in principle combine label text **with** enrolment at "
    "k ≥ 1. We give them features + 1-NN only. That is the first gap a reviewer will find.",
], 12.0)

# ================================================================ 10  divider
divider(new(), "THE PROBLEM", "Our own engine is what is holding us back",
        "Three architectures, three full runs, one consistent pattern.")

# ================================================================ 11  trajectory
s = new()
title(s, "Every redesign fixed the readout and cost us the representation")
deckline(s, "Four trained versions, mean over k = 1…16.")
picture(s, "trajectory", 0.60, 1.42, 12.13, 4.72)
banner(s, "PB-03's best-ever engine (58.6 / 48.5) still loses to the pre-engine checkpoint's "
          "plain 1-NN (61.7 / 54.6). We have been trading the part that works for the part that "
          "does not.", top=6.30)

# ================================================================ 12  ladder figure
s = new()
title(s, "Which part of the engine actually earns anything")
deckline(s, "Starting from encoder + 1-NN, adding one engine component at a time, at every k.")
picture(s, "ladder", 0.60, 1.42, 12.13, 4.72)
banner(s, "The entire native-engine gain is un-pooling. Corpus memory costs 0.24 to 0.00. The "
          "learned reranker never moves the score by more than 0.14 F1 in either direction.",
       top=6.30)

# ================================================================ 13  ladder tables
s = new()
title(s, "The same ladder, as numbers")
deckline(s, "Absolute macro F1, and the change each step makes.")
colhead(s, 0.60, 1.50, 5.90, "Ordinary")
table(s, ladder_rows("ordinary"), 0.60, 1.92, [2.35, 0.71, 0.71, 0.71, 0.71, 0.71],
      row_h=0.44, size=9.5, hi_rows=(1,))
colhead(s, 6.83, 1.50, 5.90, "Specialized / clinical")
table(s, ladder_rows("specialized_novel"), 6.83, 1.92, [2.35, 0.71, 0.71, 0.71, 0.71, 0.71],
      row_h=0.44, size=9.5, hi_rows=(1,))
write(textbox(s, 0.60, 4.25, 12.13, 1.30),
      ["Retrieval over the corpus **alone** scores 20.7 ordinary and about 10 specialized. "
       "Retrieval is the mechanism this architecture is built around, and on the deployed path "
       "it is contributing nothing."], 12.5, INK, space=4)
banner(s, "Un-pooling is a data-representation choice, not learning. Nothing we trained is "
          "carrying the result.", top=5.70)

# ================================================================ 14  decisions
s = new()
title(s, "Three decisions")
items = [
    ("1", "Do we keep trying to make the engine win?",
     "Three architectures, three runs. PB-03 finally beats its own 1-NN — but only by un-pooling, "
     "and its encoder is 4 to 8 F1 worse than the pre-engine one.",
     "→ Report the encoder + 1-NN result as the headline, and the engine as a measured negative "
     "with a named cause."),
    ("2", "Do we chase the encoder regression instead?",
     "The 4–8 F1 the encoder lost under episodic training is larger than anything the engine has "
     "ever offered us.",
     "→ Yes. This is now the load-bearing experiment: find out why episodic training degrades "
     "the representation."),
    ("3", "Do we re-train against the right control?",
     "Every run so far optimised against pooled-execution 1-NN. The ladder shows the real bar is "
     "un-pooled row 1-NN, which the engine has never once beaten.",
     "→ Change the training control before spending another 35k steps on the engine."),
]
for i, (num, head, why, rec) in enumerate(items):
    x = 0.60 + i * 4.12
    panel(s, x, 1.45, 3.85, 4.30)
    write(textbox(s, x + 0.18, 1.61, 0.50, 0.30), num, 16, BLUE, bold=True)
    write(textbox(s, x + 0.62, 1.61, 3.05, 0.80), head, 13.5, INK, bold=True, space=0)
    write(textbox(s, x + 0.20, 2.70, 3.50, 1.60), why, 12, INK, space=0)
    write(textbox(s, x + 0.20, 4.45, 3.50, 1.10), rec, 12, GREEN, space=0)
banner(s, "The strongest thing we own is a 1 M-parameter encoder that beats every baseline at "
          "every k on clinical motion, with no deployment-time training at all. Lead with that.",
       top=6.15)

# ================================================================ 15  caveats
s = new()
title(s, "What we should state plainly")
deckline(s, "Everything a reviewer will find, said first.")
colhead(s, 0.60, 1.55, 5.85, "Limits of this result")
bullets(s, 0.60, 1.97, 5.85, 3.60, [
    "**Single training seed** per HALO checkpoint; our own policy asks for three.",
    "**Zero-shot on specialized activities is weak** (8.8 against UniMTS's 17.4).",
    "**The best checkpoint is pre-engine**, and the current code cannot even load it — extending "
    "it needs a retrain.",
    "**LiMU-BERT and CrossHAR released no weights**, so we pretrained them on our corpus. That "
    "cuts both ways and we should say so.",
], 12.5)
colhead(s, 6.85, 1.55, 5.85, "Confounds we have not resolved")
bullets(s, 6.85, 1.97, 5.85, 3.60, [
    "The best checkpoint's recipe also differed — learnable front end, alias episodes, "
    "max_support 4 — so **“the engine broke the encoder” is the leading explanation, not a "
    "proven one**.",
    "We have not given the text-aligned baselines a text + enrolment readout, which is arguably "
    "part of their contribution.",
    "Un-pooling helps our engine; we have not tested whether it would help the baselines too.",
], 12.5)
banner(s, "None of these change the headline. All of them belong in the paper before a reviewer "
          "finds them.", top=6.15)

# ---------------------------------------------------------------- geometry check
SLIDE_W, SLIDE_H = 13.3330, 7.5
problems = []
for n, sl in enumerate(prs.slides, start=1):
    boxes = []
    for sh in sl.shapes:
        l, t = Emu(sh.left).inches, Emu(sh.top).inches
        r, b = l + Emu(sh.width).inches, t + Emu(sh.height).inches
        if l < -0.01 or t < -0.01 or r > SLIDE_W + 0.01 or b > SLIDE_H + 0.01:
            problems.append(f"slide {n}: {sh.shape_type} {sh.name} out of bounds "
                            f"({l:.2f},{t:.2f})-({r:.2f},{b:.2f})")
        if sh.shape_type is not None and "PICTURE" in str(sh.shape_type):
            boxes.append(("pic", l, t, r, b))
        elif sh.has_table:
            boxes.append(("table", l, t, r, b))
    banners = [(Emu(sh.top).inches, Emu(sh.top).inches + Emu(sh.height).inches)
               for sh in sl.shapes
               if sh.shape_type is not None and "AUTO_SHAPE" in str(sh.shape_type)
               and getattr(sh.fill, "type", None) == 1
               and str(sh.fill.fore_color.rgb) == CREAM]
    for kind, l, t, r, b in boxes:
        for bt, bb in banners:
            if b > bt + 0.01 and t < bb:
                problems.append(f"slide {n}: {kind} bottom {b:.2f} overlaps banner at {bt:.2f}")
print("geometry:", "OK" if not problems else "PROBLEMS")
for line in problems:
    print("  !", line)

out = ROOT / "HALO_Weekly_20260824.pptx"
prs.save(out)
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides, "
      f"{out.stat().st_size/1e6:.2f} MB)")
